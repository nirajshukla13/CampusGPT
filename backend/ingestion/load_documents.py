import os
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


def load_document(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".txt":
        return parse_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def parse_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = []

    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)

    # ✅ RETURN STRING, NOT LIST
    return "\n\n".join(pages)


def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    paragraphs = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    return "\n\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides_text = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"Slide {slide_index}:"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text)
        slides_text.append("\n".join(slide_lines))

    return "\n\n".join(slides_text)


def parse_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
